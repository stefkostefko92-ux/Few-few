#!/usr/bin/env python3
"""Generate a 9-slide presentation: Doors and Windows (Career Paths - Construction, Unit 9)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Colour palette (inspired by the orange/blue textbook cover)
ORANGE = RGBColor(0xE8, 0x8A, 0x1A)
DARK_BLUE = RGBColor(0x2A, 0x3D, 0x66)
LIGHT_BLUE = RGBColor(0xE9, 0xED, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x44, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def band(slide, top, height, color):
    shp = slide.shapes.add_shape(1, 0, top, SW, height)  # rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, left, top, width, height, lines, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def bullets(slide, left, top, width, height, items, size=18,
            color=GREY, bullet_color=ORANGE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        if isinstance(item, tuple):
            term, definition = item
            r1 = p.add_run(); r1.text = f"▪ {term}: "
            r1.font.bold = True; r1.font.size = Pt(size); r1.font.color.rgb = DARK_BLUE
            r2 = p.add_run(); r2.text = definition
            r2.font.size = Pt(size); r2.font.color.rgb = color
            r1.font.name = "Calibri"; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = f"▪ {item}"
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def slide_header(slide, number, title):
    band(slide, 0, Inches(1.15), DARK_BLUE)
    band(slide, Inches(1.15), Inches(0.12), ORANGE)
    textbox(slide, Inches(0.5), 0, Inches(11.0), Inches(1.15),
            [(title, 30, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    # number badge
    badge = slide.shapes.add_shape(9, Inches(12.0), Inches(0.28), Inches(0.6), Inches(0.6))
    badge.fill.solid(); badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background(); badge.shadow.inherit = False
    bf = badge.text_frame; bf.word_wrap = True
    p = bf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(number); r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = WHITE


# ---------------------------------------------------------------- Slide 1 (title)
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_BLUE)
band(s, Inches(2.7), Inches(2.1), ORANGE)
textbox(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(1.0),
        [("DOORS AND WINDOWS", 48, WHITE, True)], align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.8), Inches(3.85), Inches(11.7), Inches(0.7),
        [("Construction – Unit 9  |  Describing Options", 22, WHITE, False)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.0),
        [("Career Paths: Construction", 18, RGBColor(0xCC, 0xD3, 0xE5), False),
         ("Reading context: Product descriptions", 16, RGBColor(0xCC, 0xD3, 0xE5), False)],
        align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------- Slide 2 (intro)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 2, "Introduction")
bullets(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(5.4), [
    "Doors and windows are key openings in the building envelope.",
    "They control access, light, ventilation, views and security.",
    "Correct selection depends on function, location, climate and design.",
    "This unit helps you describe and compare the available options.",
    "Key skill (Function): describing options to a client or co-worker.",
], size=22)

# ---------------------------------------------------------------- Slide 3 (door parts / types overview)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 3, "Doors: Main Types")
band(s, Inches(1.4), Inches(5.9), LIGHT_BLUE)
bullets(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(5.4), [
    ("Sliding door", "moves sideways on a track."),
    ("Pocket door", "slides into a cavity inside the wall."),
    ("Bifold door", "folds in two panels to open."),
    ("Revolving door", "rotating panels; saves heated/cooled air."),
], size=19)
bullets(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(5.4), [
    ("Automatic door", "opens by sensor, no manual effort."),
    ("Pocket / fire door", "fire door resists fire and smoke spread."),
    ("Fire door", "rated to slow the spread of fire."),
    ("Frame", "the surround that holds the door in place."),
], size=19)

# ---------------------------------------------------------------- Slide 4 (windows types)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 4, "Windows: Main Types")
bullets(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(5.4), [
    ("Fixed", "does not open; provides light and view only."),
    ("Double-hung", "two sashes slide up and down."),
    ("Casement", "hinged at the side; swings outward."),
], size=19)
bullets(s, Inches(6.9), Inches(1.7), Inches(6.0), Inches(5.4), [
    ("Awning", "hinged at top; opens outward from the bottom."),
    ("Hopper", "hinged at bottom; opens inward from the top."),
    ("Sliding", "sash slides horizontally on a track."),
], size=19)

# ---------------------------------------------------------------- Slide 5 (window components)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 5, "Window Components")
bullets(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4), [
    ("Sash", "the moveable part of the window that holds the glass."),
    ("Pane", "a single sheet of glass within the sash."),
    ("Frame", "the fixed outer structure fitted into the wall opening."),
    ("Glazing", "the glass and the way it is set into the sash."),
], size=21)

# ---------------------------------------------------------------- Slide 6 (key vocabulary)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 6, "Key Vocabulary")
band(s, Inches(1.4), Inches(5.9), LIGHT_BLUE)
vocab = ["automatic door", "awning", "bifold door", "casement",
         "double-hung", "fire door", "fixed", "frame", "sash",
         "hopper", "pane", "pocket door", "revolving door", "sliding"]
col1 = vocab[:7]; col2 = vocab[7:]
bullets(s, Inches(0.9), Inches(1.8), Inches(5.8), Inches(5.2), col1, size=20)
bullets(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.2), col2, size=20)

# ---------------------------------------------------------------- Slide 7 (choosing / describing options)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 7, "Choosing the Right Option")
bullets(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.4), [
    ("Space", "sliding and pocket doors save floor space."),
    ("Energy", "revolving and automatic doors reduce air loss."),
    ("Safety", "fire doors are required on escape routes."),
    ("Ventilation", "casement and awning windows give good airflow."),
    ("View & light", "fixed and large panes maximise daylight."),
], size=20)

# ---------------------------------------------------------------- Slide 8 (useful language)
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
slide_header(s, 8, "Useful Language: Describing Options")
band(s, Inches(1.4), Inches(5.9), LIGHT_BLUE)
bullets(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.2), [
    "“One option is to install a sliding door.”",
    "“Another possibility would be a casement window.”",
    "“You could go with a fire door for the corridor.”",
    "“If you want more light, I’d recommend a fixed window.”",
    "“Alternatively, we can use a revolving door at the entrance.”",
], size=21, bullet_color=ORANGE)

# ---------------------------------------------------------------- Slide 9 (summary)
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_BLUE)
band(s, Inches(0.0), Inches(1.15), ORANGE)
textbox(s, Inches(0.6), 0, Inches(12.0), Inches(1.15),
        [("Summary", 34, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.6), [
    "Doors: sliding, pocket, bifold, revolving, automatic, fire.",
    "Windows: fixed, double-hung, casement, awning, hopper, sliding.",
    "Components: frame, sash, pane.",
    "Choose by space, energy, safety, ventilation and view.",
    "Function practised: describing options.",
], size=21, color=WHITE, bullet_color=ORANGE)
textbox(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7),
        [("Thank you!", 24, ORANGE, True)], align=PP_ALIGN.CENTER)

out = "/home/user/Few-few/Doors_and_Windows.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
